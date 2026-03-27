"""
오피스 문서 출력 기능 테스트

Docker 이미지에 설치된 openpyxl, python-docx, python-pptx를 사용하여
Excel, Word, PowerPoint 파일을 생성하고 출력 디렉토리에 저장하는 기능을 테스트합니다.
"""

import os
import sys
from pathlib import Path
from datetime import datetime

# 프로젝트 루트를 Python 경로에 추가
project_root = Path(__file__).parent.parent.parent
sys.path.insert(0, str(project_root))

from src.tools.code_execution import execute_code_in_docker


def test_excel_output():
    """Excel 파일 생성 테스트"""
    print("\n" + "="*70)
    print("테스트 1: Excel 파일 생성 (openpyxl)")
    print("="*70)
    
    # 출력 디렉토리 생성
    output_dir = Path("tests/test_output_office")
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # 테스트 코드 생성
    code_dir = Path("workspace/generated_code")
    code_dir.mkdir(parents=True, exist_ok=True)
    
    test_code = """
import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill
import os

print("📊 Excel 파일 생성 시작...")

# 샘플 데이터 생성
data = {
    '제품명': ['노트북', '마우스', '키보드', '모니터'],
    '가격': [1200000, 50000, 80000, 350000],
    '재고': [15, 120, 80, 25],
    '카테고리': ['전자기기', '주변기기', '주변기기', '전자기기']
}
df = pd.DataFrame(data)

# Excel 파일 생성
output_path = "/workspace/results/제품_목록.xlsx"
wb = Workbook()
ws = wb.active
ws.title = "제품 목록"

# 헤더 스타일
header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
header_font = Font(color="FFFFFF", bold=True)

# 데이터 쓰기
for col_idx, col_name in enumerate(df.columns, 1):
    cell = ws.cell(row=1, column=col_idx, value=col_name)
    cell.fill = header_fill
    cell.font = header_font
    cell.alignment = Alignment(horizontal='center', vertical='center')

for row_idx, row_data in enumerate(df.values, 2):
    for col_idx, value in enumerate(row_data, 1):
        cell = ws.cell(row=row_idx, column=col_idx, value=value)
        if isinstance(value, (int, float)):
            cell.alignment = Alignment(horizontal='right', vertical='center')

# 열 너비 자동 조정
for column in ws.columns:
    max_length = 0
    column_letter = column[0].column_letter
    for cell in column:
        try:
            if len(str(cell.value)) > max_length:
                max_length = len(str(cell.value))
        except:
            pass
    adjusted_width = min(max_length + 2, 50)
    ws.column_dimensions[column_letter].width = adjusted_width

wb.save(output_path)

if os.path.exists(output_path):
    print(f"✅ Excel 파일 생성 성공: {output_path}")
    print(f"   - 행 수: {len(df) + 1} (헤더 포함)")
    print(f"   - 열 수: {len(df.columns)}")
    print(f"   - 시트명: {ws.title}")
else:
    print(f"❌ Excel 파일 생성 실패: {output_path}")
"""
    
    code_file = code_dir / "test_excel_output.py"
    code_file.write_text(test_code, encoding='utf-8')
    
    print(f"\n📋 입력:")
    print(f"  코드 파일: {code_file}")
    print(f"  출력 디렉토리: {output_dir}")
    print(f"  예상 출력 파일: {output_dir}/제품_목록.xlsx")
    
    try:
        print("\n🚀 Docker에서 코드 실행 시작...")
        result = execute_code_in_docker(
            code_file=str(code_file),
            docker_image="csv-sandbox:test",
            output_directory=str(output_dir),
            timeout=60
        )
        
        print("\n✅ 실행 결과:")
        if result.success:
            print(f"  ✅ 실행 성공")
            print(f"\n  실행 출력:")
            print(result.stdout)
            
            # 출력 파일 확인
            output_file = output_dir / "제품_목록.xlsx"
            if output_file.exists():
                file_size = output_file.stat().st_size
                print(f"\n  ✅ Excel 파일 생성 확인:")
                print(f"     - 파일 경로: {output_file}")
                print(f"     - 파일 크기: {file_size:,} bytes")
                return True
            else:
                print(f"\n  ⚠️ 출력 파일이 생성되지 않았습니다: {output_file}")
                return False
        else:
            print(f"  ❌ 실행 실패")
            print(f"  에러: {result.error}")
            if result.stderr:
                print(f"  stderr: {result.stderr}")
            return False
            
    except Exception as e:
        print(f"\n❌ 오류 발생: {str(e)}")
        import traceback
        traceback.print_exc()
        return False


def test_word_output():
    """Word 문서 생성 테스트"""
    print("\n" + "="*70)
    print("테스트 2: Word 문서 생성 (python-docx)")
    print("="*70)
    
    # 출력 디렉토리 생성
    output_dir = Path("tests/test_output_office")
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # 테스트 코드 생성
    code_dir = Path("workspace/generated_code")
    code_dir.mkdir(parents=True, exist_ok=True)
    
    test_code = """
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from datetime import datetime
import os

print("📝 Word 문서 생성 시작...")

# Word 문서 생성
doc = Document()

# 제목
title = doc.add_heading('데이터 분석 보고서', 0)
title.alignment = WD_ALIGN_PARAGRAPH.CENTER

# 날짜
date_para = doc.add_paragraph(f'생성일: {datetime.now().strftime("%Y년 %m월 %d일")}')
date_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

# 목차
doc.add_heading('목차', level=1)
doc.add_paragraph('1. 요약', style='List Bullet')
doc.add_paragraph('2. 데이터 개요', style='List Bullet')
doc.add_paragraph('3. 분석 결과', style='List Bullet')
doc.add_paragraph('4. 결론', style='List Bullet')

# 본문 섹션
doc.add_heading('1. 요약', level=1)
doc.add_paragraph(
    '본 보고서는 데이터 분석 결과를 요약한 문서입니다. '
    '주요 발견 사항과 결론을 제시합니다.'
)

doc.add_heading('2. 데이터 개요', level=1)
doc.add_paragraph('분석에 사용된 데이터의 기본 정보는 다음과 같습니다:')

# 테이블 추가
table = doc.add_table(rows=5, cols=3)
table.style = 'Light Grid Accent 1'

# 헤더 행
header_cells = table.rows[0].cells
header_cells[0].text = '항목'
header_cells[1].text = '값'
header_cells[2].text = '비고'

# 데이터 행
data_rows = [
    ['샘플 수', '150', ''],
    ['특성 수', '10', ''],
    ['결측치', '0', '없음'],
    ['이상치', '3', '제외 처리']
]

for i, row_data in enumerate(data_rows, 1):
    cells = table.rows[i].cells
    for j, cell_value in enumerate(row_data):
        cells[j].text = str(cell_value)

doc.add_heading('3. 분석 결과', level=1)
doc.add_paragraph('분석 결과 주요 내용은 다음과 같습니다:')

# 번호 매기기 목록
results = [
    '평균 값은 42.5로 측정되었습니다.',
    '표준 편차는 12.3입니다.',
    '상관관계 분석 결과 유의미한 상관을 확인했습니다.'
]

for result in results:
    doc.add_paragraph(result, style='List Number')

doc.add_heading('4. 결론', level=1)
doc.add_paragraph(
    '분석 결과를 종합하면, 데이터의 품질이 우수하며 '
    '추가 분석이 가능한 상태입니다. 후속 연구를 권장합니다.'
)

# 저장
output_path = "/workspace/results/분석_보고서.docx"
doc.save(output_path)

if os.path.exists(output_path):
    file_size = os.path.getsize(output_path)
    print(f"✅ Word 문서 생성 성공: {output_path}")
    print(f"   - 파일 크기: {file_size:,} bytes")
    print(f"   - 제목 수: {len([p for p in doc.paragraphs if p.style.name.startswith('Heading')])}")
    print(f"   - 테이블 수: {len(doc.tables)}")
else:
    print(f"❌ Word 문서 생성 실패: {output_path}")
"""
    
    code_file = code_dir / "test_word_output.py"
    code_file.write_text(test_code, encoding='utf-8')
    
    print(f"\n📋 입력:")
    print(f"  코드 파일: {code_file}")
    print(f"  출력 디렉토리: {output_dir}")
    print(f"  예상 출력 파일: {output_dir}/분석_보고서.docx")
    
    try:
        print("\n🚀 Docker에서 코드 실행 시작...")
        result = execute_code_in_docker(
            code_file=str(code_file),
            docker_image="csv-sandbox:test",
            output_directory=str(output_dir),
            timeout=60
        )
        
        print("\n✅ 실행 결과:")
        if result.success:
            print(f"  ✅ 실행 성공")
            print(f"\n  실행 출력:")
            print(result.stdout)
            
            # 출력 파일 확인
            output_file = output_dir / "분석_보고서.docx"
            if output_file.exists():
                file_size = output_file.stat().st_size
                print(f"\n  ✅ Word 문서 생성 확인:")
                print(f"     - 파일 경로: {output_file}")
                print(f"     - 파일 크기: {file_size:,} bytes")
                return True
            else:
                print(f"\n  ⚠️ 출력 파일이 생성되지 않았습니다: {output_file}")
                return False
        else:
            print(f"  ❌ 실행 실패")
            print(f"  에러: {result.error}")
            if result.stderr:
                print(f"  stderr: {result.stderr}")
            return False
            
    except Exception as e:
        print(f"\n❌ 오류 발생: {str(e)}")
        import traceback
        traceback.print_exc()
        return False


def test_powerpoint_output():
    """PowerPoint 프레젠테이션 생성 테스트"""
    print("\n" + "="*70)
    print("테스트 3: PowerPoint 프레젠테이션 생성 (python-pptx)")
    print("="*70)
    
    # 출력 디렉토리 생성
    output_dir = Path("tests/test_output_office")
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # 테스트 코드 생성
    code_dir = Path("workspace/generated_code")
    code_dir.mkdir(parents=True, exist_ok=True)
    
    test_code = """
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime
import os

print("📊 PowerPoint 프레젠테이션 생성 시작...")

# 프레젠테이션 생성
prs = Presentation()

# 슬라이드 1: 제목 슬라이드
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]

title1.text = "데이터 분석 결과"
subtitle1.text = f"{datetime.now().strftime('%Y년 %m월 %d일')}"

# 슬라이드 2: 제목 및 내용
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title2 = slide2.shapes.title
content2 = slide2.placeholders[1]

title2.text = "분석 개요"
tf2 = content2.text_frame
tf2.text = "본 프레젠테이션은 데이터 분석 결과를 요약한 자료입니다."

p = tf2.add_paragraph()
p.text = "주요 내용:"
p.level = 1

p = tf2.add_paragraph()
p.text = "데이터 수집 및 전처리"
p.level = 2

p = tf2.add_paragraph()
p.text = "통계 분석 수행"
p.level = 2

p = tf2.add_paragraph()
p.text = "시각화 및 결과 해석"
p.level = 2

# 슬라이드 3: 제목 및 내용 (결과)
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title3 = slide3.shapes.title
content3 = slide3.placeholders[1]

title3.text = "주요 발견 사항"
tf3 = content3.text_frame
tf3.text = "분석 결과 다음과 같은 주요 발견 사항을 확인했습니다:"

p = tf3.add_paragraph()
p.text = "평균값: 42.5"
p.level = 1

p = tf3.add_paragraph()
p.text = "표준편차: 12.3"
p.level = 1

p = tf3.add_paragraph()
p.text = "유의미한 상관관계 확인"
p.level = 1

# 슬라이드 4: 제목 및 내용 (결론)
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title4 = slide4.shapes.title
content4 = slide4.placeholders[1]

title4.text = "결론 및 권장 사항"
tf4 = content4.text_frame
tf4.text = "분석 결과를 바탕으로 다음과 같은 결론을 도출했습니다:"

p = tf4.add_paragraph()
p.text = "데이터 품질이 우수함"
p.level = 1

p = tf4.add_paragraph()
p.text = "추가 분석이 가능한 상태"
p.level = 1

p = tf4.add_paragraph()
p.text = "후속 연구 권장"
p.level = 1

# 슬라이드 5: 제목 슬라이드 (마무리)
slide5 = prs.slides.add_slide(prs.slide_layouts[0])
title5 = slide5.shapes.title
subtitle5 = slide5.placeholders[1]

title5.text = "감사합니다"
subtitle5.text = "질문 사항은 언제든지 환영합니다"

# 저장
output_path = "/workspace/results/분석_프레젠테이션.pptx"
prs.save(output_path)

if os.path.exists(output_path):
    file_size = os.path.getsize(output_path)
    print(f"✅ PowerPoint 프레젠테이션 생성 성공: {output_path}")
    print(f"   - 파일 크기: {file_size:,} bytes")
    print(f"   - 슬라이드 수: {len(prs.slides)}")
    print(f"   - 슬라이드 레이아웃: {[s.slide_layout.name for s in prs.slides]}")
else:
    print(f"❌ PowerPoint 프레젠테이션 생성 실패: {output_path}")
"""
    
    code_file = code_dir / "test_powerpoint_output.py"
    code_file.write_text(test_code, encoding='utf-8')
    
    print(f"\n📋 입력:")
    print(f"  코드 파일: {code_file}")
    print(f"  출력 디렉토리: {output_dir}")
    print(f"  예상 출력 파일: {output_dir}/분석_프레젠테이션.pptx")
    
    try:
        print("\n🚀 Docker에서 코드 실행 시작...")
        result = execute_code_in_docker(
            code_file=str(code_file),
            docker_image="csv-sandbox:test",
            output_directory=str(output_dir),
            timeout=60
        )
        
        print("\n✅ 실행 결과:")
        if result.success:
            print(f"  ✅ 실행 성공")
            print(f"\n  실행 출력:")
            print(result.stdout)
            
            # 출력 파일 확인
            output_file = output_dir / "분석_프레젠테이션.pptx"
            if output_file.exists():
                file_size = output_file.stat().st_size
                print(f"\n  ✅ PowerPoint 프레젠테이션 생성 확인:")
                print(f"     - 파일 경로: {output_file}")
                print(f"     - 파일 크기: {file_size:,} bytes")
                return True
            else:
                print(f"\n  ⚠️ 출력 파일이 생성되지 않았습니다: {output_file}")
                return False
        else:
            print(f"  ❌ 실행 실패")
            print(f"  에러: {result.error}")
            if result.stderr:
                print(f"  stderr: {result.stderr}")
            return False
            
    except Exception as e:
        print(f"\n❌ 오류 발생: {str(e)}")
        import traceback
        traceback.print_exc()
        return False


def test_combined_output():
    """Excel, Word, PPT 동시 생성 테스트"""
    print("\n" + "="*70)
    print("테스트 4: 모든 오피스 문서 동시 생성")
    print("="*70)
    
    # 출력 디렉토리 생성
    output_dir = Path("tests/test_output_office")
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # 테스트 코드 생성
    code_dir = Path("workspace/generated_code")
    code_dir.mkdir(parents=True, exist_ok=True)
    
    test_code = """
import pandas as pd
from openpyxl import Workbook
from docx import Document
from pptx import Presentation
from datetime import datetime
import os

print("📦 모든 오피스 문서 생성 시작...")

# 1. Excel 파일 생성
print("\\n1. Excel 파일 생성 중...")
wb = Workbook()
ws = wb.active
ws.title = "요약"
data = [['항목', '값'], ['평균', 42.5], ['표준편차', 12.3], ['샘플 수', 150]]
for row in data:
    ws.append(row)
excel_path = "/workspace/results/요약.xlsx"
wb.save(excel_path)
print(f"   ✅ Excel 생성 완료: {excel_path}")

# 2. Word 문서 생성
print("\\n2. Word 문서 생성 중...")
doc = Document()
doc.add_heading('요약 보고서', 0)
doc.add_paragraph(f'생성일: {datetime.now().strftime("%Y-%m-%d")}')
doc.add_paragraph('데이터 분석 결과를 요약한 문서입니다.')
word_path = "/workspace/results/요약.docx"
doc.save(word_path)
print(f"   ✅ Word 생성 완료: {word_path}")

# 3. PowerPoint 프레젠테이션 생성
print("\\n3. PowerPoint 프레젠테이션 생성 중...")
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "요약"
slide.placeholders[1].text = "데이터 분석 결과"
ppt_path = "/workspace/results/요약.pptx"
prs.save(ppt_path)
print(f"   ✅ PowerPoint 생성 완료: {ppt_path}")

print("\\n✅ 모든 문서 생성 완료!")
print(f"\\n생성된 파일 목록:")
for file_path in [excel_path, word_path, ppt_path]:
    if os.path.exists(file_path):
        size = os.path.getsize(file_path)
        print(f"  - {os.path.basename(file_path)}: {size:,} bytes")
"""
    
    code_file = code_dir / "test_combined_output.py"
    code_file.write_text(test_code, encoding='utf-8')
    
    print(f"\n📋 입력:")
    print(f"  코드 파일: {code_file}")
    print(f"  출력 디렉토리: {output_dir}")
    
    try:
        print("\n🚀 Docker에서 코드 실행 시작...")
        result = execute_code_in_docker(
            code_file=str(code_file),
            docker_image="csv-sandbox:test",
            output_directory=str(output_dir),
            timeout=60
        )
        
        print("\n✅ 실행 결과:")
        if result.success:
            print(f"  ✅ 실행 성공")
            print(f"\n  실행 출력:")
            print(result.stdout)
            
            # 모든 출력 파일 확인
            expected_files = [
                output_dir / "요약.xlsx",
                output_dir / "요약.docx",
                output_dir / "요약.pptx"
            ]
            
            all_exist = True
            for output_file in expected_files:
                if output_file.exists():
                    file_size = output_file.stat().st_size
                    print(f"\n  ✅ {output_file.name} 생성 확인 ({file_size:,} bytes)")
                else:
                    print(f"\n  ❌ {output_file.name} 생성 실패")
                    all_exist = False
            
            return all_exist
        else:
            print(f"  ❌ 실행 실패")
            print(f"  에러: {result.error}")
            if result.stderr:
                print(f"  stderr: {result.stderr}")
            return False
            
    except Exception as e:
        print(f"\n❌ 오류 발생: {str(e)}")
        import traceback
        traceback.print_exc()
        return False


def main():
    """메인 테스트 함수"""
    print("\n" + "="*70)
    print("오피스 문서 출력 기능 테스트 시작")
    print(f"시작 시간: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    print("="*70)
    
    results = []
    
    # 테스트 실행
    print("\n" + "="*70)
    print("테스트 실행 중...")
    print("="*70)
    
    results.append(("Excel 파일 생성", test_excel_output()))
    results.append(("Word 문서 생성", test_word_output()))
    results.append(("PowerPoint 프레젠테이션 생성", test_powerpoint_output()))
    results.append(("모든 문서 동시 생성", test_combined_output()))
    
    # 결과 요약
    print("\n" + "="*70)
    print("테스트 결과 요약")
    print("="*70)
    
    for test_name, passed in results:
        status = "✅ 통과" if passed else "❌ 실패"
        print(f"  {test_name}: {status}")
    
    total = len(results)
    passed = sum(1 for _, p in results if p)
    
    print(f"\n총 {total}개 테스트 중 {passed}개 통과 ({passed/total*100:.1f}%)")
    print(f"종료 시간: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    
    return all(p for _, p in results)


if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)

